import os
import uuid
from copy import deepcopy
from datetime import date

from flask import (
    Blueprint,
    abort,
    current_app,
    flash,
    redirect,
    render_template,
    request,
    send_file,
    url_for,
)
from werkzeug.utils import secure_filename

from models import DocumentTemplate, Member, MemberAward, db
from utils import login_required, parse_date

bp = Blueprint("awards", __name__, url_prefix="/awards")


def _get_context(template, member, issued_at=None):
    return {
        "member": member,
        "today": date.today(),
        "issued_at": issued_at or date.today(),
        "template": template,
        "image_url": template.image_url,
    }


def _save_template_image(file):
    if not file or not file.filename:
        return None
    ext = os.path.splitext(secure_filename(file.filename).lower())[1]
    if ext not in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        return None
    upload_dir = os.path.join(current_app.config["UPLOAD_FOLDER"], "templates")
    os.makedirs(upload_dir, exist_ok=True)
    filename = f"{uuid.uuid4().hex}{ext}"
    file.save(os.path.join(upload_dir, filename))
    return os.path.join("templates", filename).replace("\\", "/")


def _save_pptx(file):
    if not file or not file.filename:
        return None
    ext = os.path.splitext(secure_filename(file.filename).lower())[1]
    if ext not in {".pptx"}:
        return None
    upload_dir = os.path.join(current_app.config["UPLOAD_FOLDER"], "pptx")
    os.makedirs(upload_dir, exist_ok=True)
    filename = f"{uuid.uuid4().hex}{ext}"
    file.save(os.path.join(upload_dir, filename))
    return os.path.join("pptx", filename).replace("\\", "/")


def _delete_uploaded(path):
    if not path:
        return
    try:
        os.remove(os.path.join(current_app.config["UPLOAD_FOLDER"], path))
    except OSError:
        pass


def _extract_pptx_shapes(path):
    try:
        from pptx import Presentation
    except ImportError:
        return []

    full_path = os.path.join(current_app.config["UPLOAD_FOLDER"], path)
    prs = Presentation(full_path)
    shapes = []
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = " ".join(p.text for p in shape.text_frame.paragraphs if p.text).strip()
            if text:
                shapes.append({
                    "id": str(shape.shape_id),
                    "text": text[:120],
                })
    return shapes


def _pptx_value(name, member, issued_at, today, note=None):
    if name == "full_name":
        return member.full_name or ""
    if name == "department":
        return member.department or ""
    if name == "position":
        return member.position or ""
    if name == "issued_at":
        return issued_at.strftime("%d.%m.%Y") if issued_at else ""
    if name == "today":
        return today.strftime("%d.%m.%Y")
    if name == "note":
        return note or ""
    return ""


def _apply_mapping_to_slide(slide, mapping, member, issued_at, today, note=None):
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        key = mapping.get(str(shape.shape_id))
        if not key:
            continue
        value = _pptx_value(key, member, issued_at, today, note)
        text_frame = shape.text_frame
        if text_frame.paragraphs:
            first_para = text_frame.paragraphs[0]
            if first_para.runs:
                first_para.runs[0].text = value
                for run in first_para.runs[1:]:
                    run.text = ""
            else:
                first_para.text = value


def _duplicate_slide(prs, index):
    source = prs.slides[index]
    blank_layout = source.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)
    for shape in list(new_slide.shapes):
        new_slide.shapes._spTree.remove(shape.element)
    for shape in source.shapes:
        newel = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    return new_slide


def _generate_pptx(template, members, issued_at, note=None):
    from pptx import Presentation

    full_path = os.path.join(current_app.config["UPLOAD_FOLDER"], template.pptx_path)
    prs = Presentation(full_path)
    today = date.today()
    mapping = template.shape_map
    members_list = list(members)
    if not members_list:
        raise ValueError("Нет членов для генерации")

    _apply_mapping_to_slide(prs.slides[0], mapping, members_list[0], issued_at, today, note)
    for member in members_list[1:]:
        new_slide = _duplicate_slide(prs, len(prs.slides) - 1)
        _apply_mapping_to_slide(new_slide, mapping, member, issued_at, today, note)

    import io

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


@bp.route("/")
@login_required
def index():
    templates = DocumentTemplate.query.filter_by(active=True).order_by(DocumentTemplate.order).all()
    awards = MemberAward.query.order_by(MemberAward.issued_at.desc()).limit(50).all()
    return render_template("awards/index.html", templates=templates, awards=awards)


@bp.route("/templates/add", methods=["POST"])
@login_required
def add_template():
    name = (request.form.get("name") or "").strip()
    title = (request.form.get("title") or "").strip()
    body = (request.form.get("body") or "").strip()
    image_path = _save_template_image(request.files.get("image"))
    pptx_path = _save_pptx(request.files.get("pptx"))
    has_content = body or pptx_path
    if name and title and has_content:
        template = DocumentTemplate(
            name=name,
            type=request.form.get("type", "award"),
            title=title,
            body=body,
            image_path=image_path,
            pptx_path=pptx_path,
            order=request.form.get("order", type=int) or 0,
        )
        db.session.add(template)
        db.session.commit()
        flash("Шаблон добавлен", "success")
        if pptx_path:
            return redirect(url_for("awards.template_map", id=template.id))
    else:
        flash("Заполните название, заголовок и тело шаблона или загрузите PPTX", "danger")
    return redirect(url_for("awards.index"))


@bp.route("/templates/<int:id>/edit", methods=["POST"])
@login_required
def edit_template(id):
    template = db.session.get(DocumentTemplate, id) or abort(404)
    template.name = (request.form.get("name") or "").strip()
    template.title = (request.form.get("title") or "").strip()
    template.body = (request.form.get("body") or "").strip()
    template.type = request.form.get("type", "award")
    template.order = request.form.get("order", type=int) or 0
    template.active = bool(request.form.get("active"))
    old_image_path = template.image_path
    new_image_path = _save_template_image(request.files.get("image"))
    if new_image_path:
        template.image_path = new_image_path
    old_pptx_path = template.pptx_path
    new_pptx_path = _save_pptx(request.files.get("pptx"))
    if new_pptx_path:
        template.pptx_path = new_pptx_path
        template.pptx_shape_map = None
    has_content = template.body or template.pptx_path
    if template.name and template.title and has_content:
        db.session.commit()
        if new_image_path and old_image_path:
            _delete_uploaded(old_image_path)
        if new_pptx_path and old_pptx_path:
            _delete_uploaded(old_pptx_path)
        flash("Шаблон обновлен", "success")
        if new_pptx_path:
            return redirect(url_for("awards.template_map", id=template.id))
    else:
        flash("Заполните все поля", "danger")
    return redirect(url_for("awards.index"))


@bp.route("/templates/<int:id>/delete", methods=["POST"])
@login_required
def delete_template(id):
    template = db.session.get(DocumentTemplate, id) or abort(404)
    image_path = template.image_path
    pptx_path = template.pptx_path
    db.session.delete(template)
    db.session.commit()
    _delete_uploaded(image_path)
    _delete_uploaded(pptx_path)
    flash("Шаблон удален", "success")
    return redirect(url_for("awards.index"))


@bp.route("/mass")
@login_required
def mass():
    templates = DocumentTemplate.query.filter_by(active=True).order_by(DocumentTemplate.order).all()
    members = Member.query.filter_by(status="active").order_by(Member.full_name).all()
    return render_template("awards/mass.html", templates=templates, members=members)


@bp.route("/mass/issue", methods=["POST"])
@login_required
def mass_issue():
    template_id = request.form.get("template_id", type=int)
    member_ids = [int(x) for x in request.form.getlist("member_ids") if x.isdigit()]
    issued_at = parse_date(request.form.get("issued_at")) or date.today()
    if not template_id or not member_ids:
        flash("Выберите шаблон и хотя бы одного члена", "danger")
        return redirect(url_for("awards.mass"))
    template = db.session.get(DocumentTemplate, template_id) or abort(404)
    members = Member.query.filter(Member.id.in_(member_ids)).all()
    for member in members:
        db.session.add(
            MemberAward(member_id=member.id, template_id=template.id, issued_at=issued_at)
        )
    db.session.commit()
    flash(f"Выдано {len(members)} наград", "success")
    if template.is_pptx:
        try:
            pptx_io = _generate_pptx(template, members, issued_at)
            return send_file(
                pptx_io,
                as_attachment=True,
                download_name="nagrady.pptx",
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception as e:
            current_app.logger.exception("Ошибка генерации PPTX")
            flash(f"Ошибка генерации PPTX: {e}", "danger")
            return redirect(url_for("awards.mass"))
    return redirect(
        url_for("awards.print_preview", template_id=template.id, member_ids=",".join(str(m.id) for m in members))
    )


@bp.route("/print")
@login_required
def print_preview():
    template_id = request.args.get("template_id", type=int)
    member_ids = [int(x) for x in request.args.get("member_ids", "").split(",") if x.isdigit()]
    template = db.session.get(DocumentTemplate, template_id) or abort(404)
    members = Member.query.filter(Member.id.in_(member_ids)).order_by(Member.full_name).all()
    pages = []
    for member in members:
        award = (
            MemberAward.query.filter_by(member_id=member.id, template_id=template.id)
            .order_by(MemberAward.issued_at.desc())
            .first()
        )
        issued_at = award.issued_at if award else date.today()
        pages.append(template.render(_get_context(template, member, issued_at)))
    return render_template("awards/print.html", template=template, pages=pages, members=members)


@bp.route("/templates/<int:id>/map", methods=["GET", "POST"])
@login_required
def template_map(id):
    template = db.session.get(DocumentTemplate, id) or abort(404)
    if not template.pptx_path:
        flash("Для этого шаблона не загружен PPTX", "danger")
        return redirect(url_for("awards.index"))
    if request.method == "POST":
        mapping = {}
        for shape_id in request.form.getlist("shape_id"):
            key = request.form.get(f"map_{shape_id}")
            if key:
                mapping[shape_id] = key
        template.shape_map = mapping
        db.session.commit()
        flash("Соответствие полей сохранено", "success")
        return redirect(url_for("awards.index"))
    shapes = _extract_pptx_shapes(template.pptx_path)
    current_map = template.shape_map
    return render_template("awards/map.html", template=template, shapes=shapes, current_map=current_map)


@bp.route("/<int:id>/member_add", methods=["POST"])
@login_required
def member_add_award(id):
    member = db.session.get(Member, id) or abort(404)
    template_id = request.form.get("template_id", type=int)
    issued_at = parse_date(request.form.get("issued_at")) or date.today()
    note = (request.form.get("note") or "").strip() or None
    if template_id:
        db.session.add(
            MemberAward(member_id=member.id, template_id=template_id, issued_at=issued_at, note=note)
        )
        db.session.commit()
        flash("Награда добавлена", "success")
    else:
        flash("Выберите шаблон", "danger")
    return redirect(url_for("members.detail", id=member.id))


@bp.route("/<int:id>/delete", methods=["POST"])
@login_required
def delete_award(id):
    award = db.session.get(MemberAward, id) or abort(404)
    member_id = award.member_id
    db.session.delete(award)
    db.session.commit()
    flash("Награда удалена", "success")
    return redirect(url_for("members.detail", id=member_id))
